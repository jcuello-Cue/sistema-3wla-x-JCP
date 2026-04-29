# generar_ppt_python.py — Generador PPT en Python puro (sin Node.js)

import json
import sys
import os
from datetime import datetime, date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

C = {
    "navy":      "1E3A5F",
    "blue":      "2563EB",
    "light":     "DBEAFE",
    "white":     "FFFFFF",
    "gray":      "64748B",
    "lightgray": "F1F5F9",
    "red":       "EF4444",
    "orange":    "F97316",
    "yellow":    "EAB308",
    "green":     "22C55E",
    "text":      "1E293B",
}

RESP_COLORS = {
    "Fe Grande":      "EF4444",
    "CODELCO":        "EAB308",
    "R&Q Ingeniería": "F97316",
    "Externo":        "8B5CF6",
    "Otro":           "64748B",
    "Sin asignar":    "D1D5DB",
}

def fmt(n):
    return f"{n:,.1f}".replace(",","X").replace(".",",").replace("X",".")

def add_bg(slide, color):
    from pptx.util import Inches
    bg = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(color)
    bg.line.fill.background()

def add_rect(slide, x, y, w, h, fill_color, line=False):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    if not line:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=12, bold=False, color="1E293B", align="left", italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT if align=="left" else (PP_ALIGN.CENTER if align=="center" else PP_ALIGN.RIGHT)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)

def add_table_slide(slide, headers, rows, x, y, w, col_widths=None):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    row_h  = Inches(0.32)
    table  = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), row_h * n_rows).table

    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)

    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(C["navy"])
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = rgb(C["white"])

    # Rows
    for i, row in enumerate(rows):
        bg = C["lightgray"] if i % 2 == 0 else C["white"]
        for j, cell_data in enumerate(row):
            cell = table.cell(i+1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(bg)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            if isinstance(cell_data, dict):
                run.text = str(cell_data.get("text",""))
                run.font.bold = cell_data.get("bold", False)
                run.font.color.rgb = rgb(cell_data.get("color", C["text"]))
            else:
                run.text = str(cell_data)
                run.font.color.rgb = rgb(C["text"])
            run.font.size = Pt(9)


def calc_acumulado(estado, fecha_str):
    tri      = estado["trisemanal"]
    registro = estado.get("registro", {})
    hoy      = date.fromisoformat(fecha_str)
    fechas   = tri["fechas_s1"]

    hh_ej = hh_esp = 0
    por_dia  = {}
    por_area = {}
    por_resp = {}

    for fs in fechas:
        fd = date.fromisoformat(fs)
        if fd > hoy:
            por_dia[fs] = {"esp":0, "ej":0}
            continue
        esp_d = ej_d = 0
        for a in tri["actividades"]:
            if not a["inicio"] or not a["termino"]: continue
            ini = date.fromisoformat(a["inicio"])
            ter = date.fromisoformat(a["termino"])
            if ini <= fd <= ter:
                esp_d += a["hh_dia"]
                area = a["area"]
                if area not in por_area: por_area[area] = {"esp":0,"ej":0}
                por_area[area]["esp"] += a["hh_dia"]

        reg_d = registro.get(fs, {})
        for corr, act in reg_d.items():
            a_base = next((a for a in tri["actividades"] if str(a["corr"])==corr), None)
            if not a_base: continue
            hh = (act.get("cant_ejecutada",0) or 0) * a_base["rendimiento"]
            ej_d += hh
            area = a_base["area"]
            if area not in por_area: por_area[area] = {"esp":0,"ej":0}
            por_area[area]["ej"] += hh
            ini = date.fromisoformat(a_base["inicio"])
            ter = date.fromisoformat(a_base["termino"])
            hh_esp_act = a_base["hh_dia"] if ini <= fd <= ter else 0
            deficit = hh_esp_act - hh
            if deficit > 0.5:
                resp = act.get("responsable","Sin asignar") or "Sin asignar"
                por_resp[resp] = por_resp.get(resp, 0) + deficit

        por_dia[fs] = {"esp": round(esp_d,1), "ej": round(ej_d,1)}
        hh_esp += esp_d
        hh_ej  += ej_d

    return {
        "hh_ej":   round(hh_ej,1),
        "hh_esp":  round(hh_esp,1),
        "hh_meta": tri["hh_totales_s1"],
        "por_dia": por_dia, "por_area": por_area, "por_resp": por_resp
    }


def get_no_ejecutadas(estado, fecha_str):
    tri      = estado["trisemanal"]
    registro = estado.get("registro", {})
    hoy      = date.fromisoformat(fecha_str)
    fechas   = [date.fromisoformat(d) for d in tri["fechas_s1"]]
    rows = []
    for a in tri["actividades"]:
        if not a["inicio"] or not a["termino"]: continue
        ini = date.fromisoformat(a["inicio"])
        ter = date.fromisoformat(a["termino"])
        hh_esp = hh_ej = 0
        for fd in fechas:
            if fd > hoy: continue
            if ini <= fd <= ter: hh_esp += a["hh_dia"]
            reg = registro.get(fd.isoformat(),{}).get(str(a["corr"]))
            if reg: hh_ej += (reg.get("cant_ejecutada",0) or 0) * a["rendimiento"]
        deficit = round(hh_esp - hh_ej, 1)
        if deficit > 0.5:
            resp = causa = ""
            for fd in sorted(fechas, reverse=True):
                if fd > hoy: continue
                reg = registro.get(fd.isoformat(),{}).get(str(a["corr"]))
                if reg and reg.get("responsable"):
                    resp  = reg["responsable"]
                    causa = reg.get("causa","")
                    break
            rows.append({"nombre":a["nombre"],"area":a["area"],"deficit":deficit,"resp":resp,"causa":causa})
    return sorted(rows, key=lambda x: -x["deficit"])


def generar_ppt(estado_path, fecha_str, output_path):
    with open(estado_path, "r", encoding="utf-8") as f:
        estado = json.load(f)

    tri   = estado["trisemanal"]
    acu   = calc_acumulado(estado, fecha_str)
    no_ej = get_no_ejecutadas(estado, fecha_str)
    pct   = round(acu["hh_ej"] / acu["hh_meta"] * 100, 1) if acu["hh_meta"] else 0
    fecha_label = date.fromisoformat(fecha_str).strftime("%d/%m/%Y")
    fechas_s1   = tri["fechas_s1"]

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]  # blank

    # ── SLIDE 1: PORTADA ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, C["navy"])
    add_rect(s, 0, 4.0, 10, 1.625, C["blue"])
    add_text(s, "RESUMEN AVANCE S1 — PROGRAMA 3WLA", 0.4, 0.3, 9, 0.6, size=12, color="CADCFC")
    add_text(s, "LM5 Río Blanco", 0.4, 0.85, 9, 1.1, size=42, bold=True, color=C["white"])
    add_text(s, "División El Teniente — CODELCO", 0.4, 1.95, 9, 0.6, size=18, color="93C5FD")
    add_text(s, f"Trisemanal N° {tri['num_trisemanal']}   |   S1: {fechas_s1[0]} → {fechas_s1[-1]}", 0.4, 4.1, 7, 0.5, size=13, color=C["white"])
    add_text(s, f"Emitido: {fecha_label}", 0.4, 4.65, 7, 0.4, size=11, color="CADCFC")
    add_text(s, "R&Q Ingeniería", 7.5, 4.2, 2.1, 0.9, size=12, color=C["white"], align="right")

    # ── SLIDE 2: KPIs + Gráfico barras ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, C["white"])
    add_rect(s, 0, 0, 10, 0.75, C["navy"])
    add_text(s, "RESUMEN EJECUTIVO — AVANCE ACUMULADO S1", 0.3, 0, 9.4, 0.75, size=14, bold=True, color=C["white"])

    cards = [
        ("HH Ejecutadas",   fmt(acu["hh_ej"]),                  f"{pct}% de la meta",                C["blue"]),
        ("HH Esperadas",    fmt(acu["hh_esp"]),                  f"{round(acu['hh_esp']/acu['hh_meta']*100,1)}% planificado", C["gray"]),
        ("HH Déficit",      fmt(acu["hh_esp"]-acu["hh_ej"]),    "por ejecutar",                      C["red"]),
        ("Meta S1",         fmt(acu["hh_meta"]),                  f"7 días S1",                       C["navy"]),
    ]
    for i, (label, val, sub, color) in enumerate(cards):
        x = 0.3 + i*2.4
        add_rect(s, x, 0.9, 2.2, 1.8, C["lightgray"])
        add_rect(s, x, 0.9, 0.06, 1.8, color)
        add_text(s, label, x+0.15, 0.95, 2.0, 0.35, size=10, color=C["gray"])
        add_text(s, val,   x+0.15, 1.28, 2.0, 0.75, size=26, bold=True, color=C["text"])
        add_text(s, sub,   x+0.15, 2.0,  2.0, 0.35, size=10, color=color)

    # Barra progreso
    add_text(s, "Avance S1", 0.3, 2.88, 3, 0.35, size=11, color=C["gray"])
    add_text(s, f"{pct}%", 8.5, 2.88, 1.2, 0.35, size=11, bold=True, color=C["blue"], align="right")
    add_rect(s, 0.3, 3.28, 9.4, 0.22, "E2E8F0")
    add_rect(s, 0.3, 3.28, max(0.05, 9.4*pct/100), 0.22, C["blue"])

    # Gráfico
    dias_lbl = [date.fromisoformat(f).strftime("%d/%m") for f in fechas_s1]
    esp_vals = [acu["por_dia"].get(f,{}).get("esp",0) for f in fechas_s1]
    ej_vals  = [acu["por_dia"].get(f,{}).get("ej",0)  for f in fechas_s1]

    cd = ChartData()
    cd.categories = dias_lbl
    cd.add_series("Esperado",  esp_vals)
    cd.add_series("Ejecutado", ej_vals)
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.3), Inches(3.6), Inches(9.4), Inches(1.8), cd).chart
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = rgb("BFDBFE")
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = rgb(C["blue"])
    chart.has_legend = True
    chart.legend.position = 1  # top

    # ── SLIDE 3: AVANCE POR TRAMO ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, C["white"])
    add_rect(s, 0, 0, 10, 0.75, C["navy"])
    add_text(s, "AVANCE POR TRAMO / ÁREA", 0.3, 0, 9.4, 0.75, size=14, bold=True, color=C["white"])

    areas = sorted(acu["por_area"].items(), key=lambda x: -x[1]["esp"])
    area_lbl = [a[0] for a in areas]
    area_esp = [round(a[1]["esp"],1) for a in areas]
    area_ej  = [round(a[1]["ej"],1)  for a in areas]

    cd2 = ChartData()
    cd2.categories = area_lbl
    cd2.add_series("Esperado",  area_esp)
    cd2.add_series("Ejecutado", area_ej)
    chart2 = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.3), Inches(0.9), Inches(5.8), Inches(4.5), cd2).chart
    chart2.series[0].format.fill.solid()
    chart2.series[0].format.fill.fore_color.rgb = rgb("BFDBFE")
    chart2.series[1].format.fill.solid()
    chart2.series[1].format.fill.fore_color.rgb = rgb(C["blue"])
    chart2.has_legend = True

    headers3 = ["Tramo/Área", "Esperado", "Ejecutado", "% Cumpl."]
    rows3 = []
    for area, v in areas:
        p = round(v["ej"]/v["esp"]*100) if v["esp"] > 0 else 0
        pc = "16A34A" if p>=80 else ("D97706" if p>=50 else "DC2626")
        rows3.append([area, fmt(v["esp"]), fmt(v["ej"]),
                      {"text":f"{p}%","bold":True,"color":pc}])
    add_table_slide(s, headers3, rows3, 6.3, 1.0, 3.4, col_widths=[1.5,0.7,0.7,0.5])

    # ── SLIDE 4: RESPONSABILIDADES ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, C["white"])
    add_rect(s, 0, 0, 10, 0.75, C["navy"])
    add_text(s, "HH NO EJECUTADAS POR RESPONSABLE", 0.3, 0, 9.4, 0.75, size=14, bold=True, color=C["white"])

    resp_entries = sorted(acu["por_resp"].items(), key=lambda x: -x[1])
    total_resp   = sum(v for _,v in resp_entries)

    if resp_entries:
        resp_lbl  = [r[0] for r in resp_entries]
        resp_vals = [round(r[1],1) for r in resp_entries]

        cd3 = ChartData()
        cd3.categories = resp_lbl
        cd3.add_series("HH no ejecutadas", resp_vals)
        chart3 = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.3), Inches(0.9), Inches(5.8), Inches(4.5), cd3).chart
        chart3.has_legend = False
        for i, (resp, _) in enumerate(resp_entries):
            chart3.series[0].points[i].format.fill.solid()
            chart3.series[0].points[i].format.fill.fore_color.rgb = rgb(RESP_COLORS.get(resp, C["gray"]))

        headers4 = ["Responsable", "HH no ejec.", "% del total"]
        rows4 = []
        for resp, v in resp_entries:
            p = round(v/total_resp*100,1) if total_resp > 0 else 0
            rows4.append([resp, fmt(v), f"{p}%"])
        rows4.append([{"text":"TOTAL","bold":True,"color":C["text"]}, {"text":fmt(total_resp),"bold":True,"color":C["text"]}, "100%"])
        add_table_slide(s, headers4, rows4, 6.3, 1.0, 3.4, col_widths=[1.7,0.9,0.8])

    # ── SLIDE 5: TOP ACTIVIDADES CON DÉFICIT ──────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, C["white"])
    add_rect(s, 0, 0, 10, 0.75, C["navy"])
    top = no_ej[:12]
    add_text(s, f"ACTIVIDADES CON DÉFICIT ACUMULADO (top {len(top)})", 0.3, 0, 9.4, 0.75, size=14, bold=True, color=C["white"])

    headers5 = ["Actividad","Área","HH Déficit","Responsable"]
    rows5 = []
    for a in top:
        nombre = a["nombre"][:45]+"…" if len(a["nombre"])>45 else a["nombre"]
        rows5.append([nombre, a["area"], {"text":fmt(a["deficit"]),"bold":True,"color":C["red"]}, a["resp"]])
    add_table_slide(s, headers5, rows5, 0.3, 0.85, 9.4, col_widths=[4.8,1.4,1.4,1.8])

    # ── SLIDE 6: CIERRE ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_bg(s, C["navy"])
    add_rect(s, 0, 3.9, 10, 1.725, C["blue"])

    add_text(s, f"{pct}%", 0.5, 0.2, 5, 2.8, size=88, bold=True, color=C["white"], align="center")
    add_text(s, "Avance S1 a la fecha", 0.5, 3.0, 5, 0.55, size=16, color="93C5FD", align="center")

    add_text(s, f"{fmt(acu['hh_ej'])} HH ejecutadas", 5.4, 0.7, 4.2, 0.9, size=22, bold=True, color=C["white"])
    add_text(s, f"de {fmt(acu['hh_meta'])} HH comprometidas en S1", 5.4, 1.5, 4.2, 0.6, size=14, color="CADCFC")
    add_text(s, f"Esperado a la fecha:  {fmt(acu['hh_esp'])} HH", 5.4, 2.3, 4.2, 0.5, size=13, color="93C5FD")
    add_text(s, f"Déficit acumulado:    {fmt(acu['hh_esp']-acu['hh_ej'])} HH", 5.4, 2.85, 4.2, 0.5, size=13, color="FCA5A5")

    add_text(s, "R&Q Ingeniería — LM5 Río Blanco", 0.4, 3.95, 6, 0.5, size=13, color=C["white"])
    add_text(s, f"Emitido: {fecha_label}", 0.4, 4.5, 9, 0.4, size=11, color="CADCFC")

    prs.save(output_path)
    print("OK:" + output_path)


if __name__ == "__main__":
    if len(sys.argv) != 4:
        print("Uso: python generar_ppt_python.py estado_3wla.json YYYY-MM-DD output.pptx")
        sys.exit(1)
    generar_ppt(sys.argv[1], sys.argv[2], sys.argv[3])
